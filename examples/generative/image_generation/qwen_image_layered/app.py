import os
import random
import tempfile

import gradio as gr
import numpy as np
import torch
from diffusers import QwenImageLayeredPipeline
from PIL import Image
from pptx import Presentation

BASE_DIR = os.path.dirname(__file__)
ASSETS_DIR = os.path.join(BASE_DIR, "assets")
MAX_SEED = np.iinfo(np.int32).max

device = "cuda" if torch.cuda.is_available() else "cpu"
dtype = torch.float16 if device == "cuda" else torch.float32
device_map = os.environ.get("QWEN_LAYERED_DEVICE_MAP", "balanced")

if device == "cuda" and device_map not in ("none", ""):
    pipeline = QwenImageLayeredPipeline.from_pretrained(
        "Qwen/Qwen-Image-Layered",
        torch_dtype=dtype,
        device_map=device_map,
    )
else:
    pipeline = QwenImageLayeredPipeline.from_pretrained(
        "Qwen/Qwen-Image-Layered",
        torch_dtype=dtype,
    )
    pipeline = pipeline.to(device, dtype)
if device == "cuda":
    pipeline.enable_attention_slicing()
pipeline.set_progress_bar_config(disable=None)


def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        return int(inch * 914400)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(
            img_path,
            left,
            top,
            width=px_to_emu(img_width_px),
            height=px_to_emu(img_height_px),
        )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name


def export_gallery(images):
    images = [item[0] for item in images]
    return imagelist_to_pptx(images)


def infer(
    input_image,
    seed=777,
    randomize_seed=False,
    prompt=None,
    neg_prompt=" ",
    true_guidance_scale=4.0,
    num_inference_steps=50,
    layer=4,
    cfg_norm=True,
    use_en_prompt=True,
):
    if input_image is None:
        raise gr.Error("Please upload an image to decompose.")

    if randomize_seed:
        seed = random.randint(0, MAX_SEED)

    if isinstance(input_image, list):
        input_image = input_image[0]

    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError(f"Unsupported input_image type: {type(input_image)}")

    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=device).manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,
        "cfg_normalize": cfg_norm,
        "use_en_prompt": use_en_prompt,
    }

    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]

    return output_images


examples = [
    os.path.join(ASSETS_DIR, "test_images", f"{idx}.png") for idx in range(1, 14)
]

with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.Markdown("## Qwen-Image-Layered")
        with gr.Row():
            with gr.Column():
                input_image = gr.Image(label="Input Image", image_mode="RGBA")

            with gr.Column():
                seed = gr.Slider(
                    label="Seed",
                    minimum=0,
                    maximum=MAX_SEED,
                    step=1,
                    value=0,
                )

                randomize_seed = gr.Checkbox(label="Randomize seed", value=True)

                prompt = gr.Textbox(
                    label="Prompt (Optional)",
                    placeholder="Optional prompt to guide decomposition",
                    value="",
                    lines=2,
                )
                neg_prompt = gr.Textbox(
                    label="Negative Prompt (Optional)",
                    placeholder="Optional negative prompt",
                    value=" ",
                    lines=2,
                )

                with gr.Row():
                    true_guidance_scale = gr.Slider(
                        label="True guidance scale",
                        minimum=1.0,
                        maximum=10.0,
                        step=0.1,
                        value=4.0,
                    )

                    num_inference_steps = gr.Slider(
                        label="Number of inference steps",
                        minimum=1,
                        maximum=50,
                        step=1,
                        value=50,
                    )

                    layer = gr.Slider(
                        label="Layers",
                        minimum=2,
                        maximum=10,
                        step=1,
                        value=4,
                    )

                with gr.Row():
                    cfg_norm = gr.Checkbox(
                        label="Whether enable CFG normalization", value=True
                    )
                    use_en_prompt = gr.Checkbox(
                        label="Automatic caption language if no prompt provided",
                        value=True,
                    )

                with gr.Row():
                    run_button = gr.Button("Decompose!", variant="primary")

        gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
        export_btn = gr.Button("Export as PPTX")
        export_file = gr.File(label="Download PPTX")
        export_btn.click(fn=export_gallery, inputs=gallery, outputs=export_file)

    gr.Examples(
        examples=examples,
        inputs=[input_image],
        outputs=[gallery],
        fn=infer,
        examples_per_page=14,
        cache_examples=False,
        run_on_click=True,
    )

    run_button.click(
        fn=infer,
        inputs=[
            input_image,
            seed,
            randomize_seed,
            prompt,
            neg_prompt,
            true_guidance_scale,
            num_inference_steps,
            layer,
            cfg_norm,
            use_en_prompt,
        ],
        outputs=gallery,
    )

server_name = os.environ.get("GRADIO_SERVER_NAME", "0.0.0.0")
server_port = int(os.environ.get("GRADIO_SERVER_PORT", "7869"))
demo.launch(server_name=server_name, server_port=server_port)
