import argparse
import os
import random

import numpy as np
import torch
from diffusers import QwenImageLayeredPipeline
from PIL import Image
from pptx import Presentation

MAX_SEED = np.iinfo(np.int32).max


def parse_args():
    parser = argparse.ArgumentParser(description="Decompose an image into RGBA layers.")
    parser.add_argument("--image", required=True, help="Path to input image.")
    parser.add_argument("--output-dir", default="outputs", help="Directory to save layers.")
    parser.add_argument("--layers", type=int, default=4, help="Number of layers.")
    parser.add_argument("--seed", type=int, default=777, help="Random seed.")
    parser.add_argument("--randomize-seed", action="store_true", help="Randomize seed.")
    parser.add_argument("--prompt", default=None, help="Optional prompt to guide decomposition.")
    parser.add_argument("--negative-prompt", default=" ", help="Negative prompt.")
    parser.add_argument("--true-guidance-scale", type=float, default=4.0, help="CFG scale.")
    parser.add_argument("--num-inference-steps", type=int, default=50, help="Inference steps.")
    parser.add_argument("--cfg-norm", action="store_true", help="Enable CFG normalization.")
    parser.add_argument(
        "--no-cfg-norm", dest="cfg_norm", action="store_false", help="Disable CFG normalization."
    )
    parser.set_defaults(cfg_norm=True)
    parser.add_argument(
        "--use-en-prompt",
        action="store_true",
        help="Auto-caption in English if prompt is missing.",
    )
    parser.add_argument(
        "--no-use-en-prompt",
        dest="use_en_prompt",
        action="store_false",
        help="Auto-caption in Chinese if prompt is missing.",
    )
    parser.set_defaults(use_en_prompt=True)
    parser.add_argument(
        "--resolution",
        type=int,
        default=640,
        help="Resolution bucket (640 or 1024).",
    )
    parser.add_argument(
        "--pptx",
        default=None,
        help="Optional path to write PPTX with stacked layers.",
    )
    parser.add_argument(
        "--device-map",
        default="balanced",
        choices=["balanced", "cuda", "none"],
        help="Device placement for the pipeline.",
    )
    parser.add_argument(
        "--sequential-cpu-offload",
        action="store_true",
        help="Enable sequential CPU offload for lower VRAM usage.",
    )
    return parser.parse_args()


def imagelist_to_pptx(img_files, output_path):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        return int(inch * 914400)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(
            img_path,
            left,
            top,
            width=px_to_emu(img_width_px),
            height=px_to_emu(img_height_px),
        )

    prs.save(output_path)


def main():
    args = parse_args()
    if args.randomize_seed:
        args.seed = random.randint(0, MAX_SEED)

    device = "cuda" if torch.cuda.is_available() else "cpu"
    dtype = torch.float16 if device == "cuda" else torch.float32

    device_map = args.device_map if device == "cuda" and args.device_map != "none" else None
    if device_map:
        pipeline = QwenImageLayeredPipeline.from_pretrained(
            "Qwen/Qwen-Image-Layered",
            torch_dtype=dtype,
            device_map=device_map,
        )
    else:
        pipeline = QwenImageLayeredPipeline.from_pretrained(
            "Qwen/Qwen-Image-Layered",
            torch_dtype=dtype,
        )
        pipeline = pipeline.to(device, dtype)
    if device == "cuda":
        pipeline.enable_attention_slicing()
        if args.sequential_cpu_offload:
            pipeline.enable_sequential_cpu_offload()

    pil_image = Image.open(args.image).convert("RGB").convert("RGBA")
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=device).manual_seed(args.seed),
        "true_cfg_scale": args.true_guidance_scale,
        "prompt": args.prompt,
        "negative_prompt": args.negative_prompt,
        "num_inference_steps": args.num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": args.layers,
        "resolution": args.resolution,
        "cfg_normalize": args.cfg_norm,
        "use_en_prompt": args.use_en_prompt,
    }

    os.makedirs(args.output_dir, exist_ok=True)
    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]

    image_paths = []
    for idx, image in enumerate(output_images):
        output_path = os.path.join(args.output_dir, f"layer_{idx}.png")
        image.save(output_path)
        image_paths.append(output_path)

    if args.pptx:
        imagelist_to_pptx(image_paths, args.pptx)

    print(f"Saved {len(image_paths)} layers to {args.output_dir}")
    if args.pptx:
        print(f"Saved PPTX to {args.pptx}")


if __name__ == "__main__":
    main()
